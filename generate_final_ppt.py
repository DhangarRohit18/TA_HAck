from pptx import Presentation
from pptx.util import Pt

def fill_ppt():
    prs = Presentation("TESSERACT'26_PPT_TEMPLATE.pptx")
    
    # Slide 1: Title
    slide1 = prs.slides[0]
    for shape in slide1.shapes:
        if not shape.has_text_frame: continue
        if "RETHINKING" in shape.text.upper():
            pass # Keep original 
        elif "TRACK" in shape.text.upper():
            shape.text_frame.text = "Track: AI/ML"
        elif "PROBLEM" in shape.text.upper() and "ID" in shape.text.upper():
            shape.text_frame.text = "Problem Statement ID: 2"
        elif "TEAM NAME" in shape.text.upper():
            shape.text_frame.text = "Team Name: Akatsuki"
        elif "MEMBER" in shape.text.upper():
            shape.text_frame.text = "Team Members: Rohit Dhangar, Sham Patil, Harsh Shinde, Prince Vallecha, Aditya Bora"

    # Slide 2: Proposed Solution
    slide2 = prs.slides[1]
    for shape in slide2.shapes:
        if not shape.has_text_frame: continue
        if "PROPOSED SOLUTION" in shape.text.upper(): continue
        
        if "NEUROOPS" in shape.text.upper() or "IDEA/SOLUTION" in shape.text.upper() or "PROTOTYPE" in shape.text.upper():
            tf = shape.text_frame
            tf.clear()
            p = tf.add_paragraph()
            p.text = "NEUROOPS AI: Security-First Autonomous IT System"
            p.font.bold = True
            
            bullets = [
                "Protects infrastructure from development to runtime (DevSecOps + AI Defense).",
                "Scans code for SQL Injection, secrets, and unsafe APIs before deployment.",
                "ML-powered anomaly detection for real-time DDoS and Brute-force threats.",
                "Autonomous action engine: Blocks IPs, rate limits, and isolates services instantly.",
                "IT Ops integration: Self-healing and auto-scaling to maintain SLA zero-downtime.",
                "Explainable AI: Dashboard provides clear rationale for all security maneuvers."
            ]
            for b in bullets:
                p_b = tf.add_paragraph()
                p_b.text = b
                p_b.font.size = Pt(14)
                p_b.level = 0

    # Slide 3: Technical Approach
    slide3 = prs.slides[2]
    # In the provided screenshot, Technical Approach has 2 slides (slide 3 and 5 are both labeled Technical Approach).
    # I'll fill slide 3 as the main tech stack.
    for shape in slide3.shapes:
        if not shape.has_text_frame: continue
        if "TECHNICAL" in shape.text.upper(): continue
        
        if "TECHNOLOGIES" in shape.text.upper():
            tf = shape.text_frame
            tf.clear()
            p = tf.add_paragraph()
            p.text = "Tech Stack: Python (FastAPI), Scikit-learn, PyTorch, Docker, Kubernetes, React, Prometheus, Grafana."
            p.font.size = Pt(16)

    # Slide 4: Feasibility
    slide4 = prs.slides[3]
    for shape in slide4.shapes:
        if not shape.has_text_frame: continue
        if "FEASIBILITY" in shape.text.upper(): continue
        if "ANALYSIS" in shape.text.upper():
            tf = shape.text_frame
            tf.clear()
            p = tf.add_paragraph()
            p.text = "High Feasibility: Uses industry-standard tools (Prometheus, FastAPI) with AI layers for automation."
            p_c = tf.add_paragraph()
            p_c.text = "Challenges: Minimizing false positives in anomaly detection; ensuring zero-latency in response."
            p_s = tf.add_paragraph()
            p_s.text = "Strategy: Reinforcement learning for continuous tuning from past incidents."

    prs.save("Final_NEUROOPS_PPT.pptx")
    print("Final presentation saved as Final_NEUROOPS_PPT.pptx")

if __name__ == "__main__":
    fill_ppt()
