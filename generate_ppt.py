from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()

    # Define a custom color (the dark blue from the template footer)
    footer_blue = RGBColor(0, 32, 96)

    # 1. Title Slide
    slide_layout = prs.slide_layouts[6] # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Title TESSERACT '26
    txBox = slide.shapes.add_textbox(Inches(3), Inches(0.5), Inches(4), Inches(1))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "TESSERACT '26"
    p.font.bold = True
    p.font.size = Pt(44)
    p.font.name = "Stencil" # Trying to match the retro look in the image
    p.alignment = PP_ALIGN.CENTER

    # "TITLE PAGE"
    txBox2 = slide.shapes.add_textbox(Inches(3), Inches(1.5), Inches(4), Inches(0.5))
    tf2 = txBox2.text_frame
    p2 = tf2.add_paragraph()
    p2.text = "TITLE PAGE"
    p2.font.bold = True
    p2.font.size = Pt(32)
    p2.alignment = PP_ALIGN.CENTER

    # Project Title
    txBox3 = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(9), Inches(1))
    tf3 = txBox3.text_frame
    p3 = tf3.add_paragraph()
    p3.text = "Problem Statement Title - Autonomous IT: Rethinking the Future of Self-Managing Systems"
    p3.font.bold = True
    p3.font.size = Pt(24)
    p3.alignment = PP_ALIGN.CENTER

    # Details
    txBox4 = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(8), Inches(2))
    tf4 = txBox4.text_frame
    details = [
        "Track – AI/ML",
        "Problem Statement ID - 2",
        "Team Name- Akatsuki",
        "Team Members- Rohit Dhangar, Sham Patil, Harsh Shinde, Prince Vallecha, Aditya Bora"
    ]
    for detail in details:
        p4 = tf4.add_paragraph()
        p4.text = "• " + detail
        p4.font.size = Pt(18)
        p4.space_after = Pt(10)

    # Footer
    footer_rect = slide.shapes.add_shape(6, 0, Inches(7.1), Inches(10), Inches(0.4)) # ShapeType.RECTANGLE
    footer_rect.fill.solid()
    footer_rect.fill.fore_color.rgb = footer_blue
    footer_rect.line.fill.background()

    # 2. NEUROOPS / Proposed Solution
    slide = prs.slides.add_slide(slide_layout)

    # Header
    txBox_h = slide.shapes.add_textbox(Inches(3), Inches(0.5), Inches(4), Inches(0.5))
    tf_h = txBox_h.text_frame
    p_h = tf_h.add_paragraph()
    p_h.text = "NEUROOPS"
    p_h.font.bold = True
    p_h.font.size = Pt(40)
    p_h.alignment = PP_ALIGN.CENTER

    # Proposed Solution Title
    txBox_s = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(9), Inches(0.5))
    tf_s = txBox_s.text_frame
    p_s = tf_s.add_paragraph()
    p_s.text = "Proposed Solution (Describe your Idea/Solution/Prototype)"
    p_s.font.bold = True
    p_s.font.underline = True
    p_s.font.size = Pt(20)

    # Bullets
    txBox_b = slide.shapes.add_textbox(Inches(0.8), Inches(2.8), Inches(8.5), Inches(4))
    tf_b = txBox_b.text_frame
    tf_b.word_wrap = True
    bullets = [
        "NEUROOPS AI is a security-first autonomous IT system that protects infrastructure from development to runtime",
        "Detects vulnerabilities before deployment (SQL injection, secrets, unsafe code)",
        "Monitors system behavior using ML to identify real-time threats (DDoS, brute-force, abnormal traffic)",
        "Automatically responds with zero human intervention (block IPs, rate limit, isolate services)",
        "Ensures self-healing, auto-scaling, and SLA maintenance",
        "Provides explainable AI decisions for transparency",
        "Combines DevSecOps + Runtime Defense inspired by a human immune system"
    ]
    for bullet in bullets:
        p_b = tf_b.add_paragraph()
        p_b.text = "• " + bullet
        p_b.font.size = Pt(16)
        p_b.space_after = Pt(8)

    # Footer
    footer_rect2 = slide.shapes.add_shape(6, 0, Inches(7.1), Inches(10), Inches(0.4))
    footer_rect2.fill.solid()
    footer_rect2.fill.fore_color.rgb = footer_blue
    footer_rect2.line.fill.background()

    # 3. Technical Approach
    slide = prs.slides.add_slide(slide_layout)

    # Header
    txBox_h3 = slide.shapes.add_textbox(Inches(3), Inches(0.5), Inches(4), Inches(0.5))
    tf_h3 = txBox_h3.text_frame
    p_h3 = tf_h3.add_paragraph()
    p_h3.text = "TECHNICAL APPROACH"
    p_h3.font.bold = True
    p_h3.font.size = Pt(36)
    p_h3.alignment = PP_ALIGN.CENTER

    # Technologies Involved
    txBox_t = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(9), Inches(1))
    tf_t = txBox_t.text_frame
    p_t = tf_t.add_paragraph()
    p_t.text = "Technologies Involved"
    p_t.font.bold = True
    p_t.font.size = Pt(20)

    # Bullet Tech
    txBox_bt = slide.shapes.add_textbox(Inches(0.8), Inches(2.3), Inches(8.5), Inches(2.5))
    tf_bt = txBox_bt.text_frame
    techs = [
        "Backend: Python (FastAPI)",
        "AI/ML: Scikit-learn (Anomaly Detection), PyTorch",
        "DevOps: Docker, Kubernetes",
        "Monitoring: Prometheus, Grafana",
        "Frontend: React (Dashboard UI)",
        "AI Agents / Logic: LangChain / Rule-based Engine",
        "Cloud: AWS / GCP"
    ]
    for tech in techs:
        p_bt = tf_bt.add_paragraph()
        p_bt.text = "  - " + tech
        p_bt.font.size = Pt(14)

    # Implementation Methodology
    txBox_im = slide.shapes.add_textbox(Inches(0.5), Inches(5.0), Inches(9), Inches(0.5))
    tf_im = txBox_im.text_frame
    p_im = tf_im.add_paragraph()
    p_im.text = "Implementation Methodology"
    p_im.font.bold = True
    p_im.font.size = Pt(20)

    # Bullets methodology
    txBox_bm = slide.shapes.add_textbox(Inches(0.8), Inches(5.5), Inches(8.5), Inches(1.5))
    tf_bm = txBox_bm.text_frame
    methodologies = [
        "Development workflow and architecture",
        "Process diagrams / flowcharts",
        "Screenshots or working prototype visuals"
    ]
    for method in methodologies:
        p_bm = tf_bm.add_paragraph()
        p_bm.text = "o " + method
        p_bm.font.size = Pt(16)

    # Footer
    footer_rect3 = slide.shapes.add_shape(6, 0, Inches(7.1), Inches(10), Inches(0.4))
    footer_rect3.fill.solid()
    footer_rect3.fill.fore_color.rgb = footer_blue
    footer_rect3.line.fill.background()

    # 4. Feasibility and Viability
    slide = prs.slides.add_slide(slide_layout)

    # Header
    txBox_h4 = slide.shapes.add_textbox(Inches(2), Inches(0.5), Inches(6), Inches(0.5))
    tf_h4 = txBox_h4.text_frame
    p_h4 = tf_h4.add_paragraph()
    p_h4.text = "FEASIBILITY AND VIABILITY"
    p_h4.font.bold = True
    p_h4.font.size = Pt(36)
    p_h4.alignment = PP_ALIGN.CENTER

    # Bullets
    txBox_f = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(8.5), Inches(3))
    tf_f = txBox_f.text_frame
    fv_bullets = [
        "Analysis of the feasibility of the idea",
        "Potential challenges and risks",
        "Strategies for overcoming these challenges"
    ]
    for bullet in fv_bullets:
        p_f = tf_f.add_paragraph()
        p_f.text = "• " + bullet
        p_f.font.size = Pt(22)
        p_f.space_after = Pt(20)

    # Footer
    footer_rect4 = slide.shapes.add_shape(6, 0, Inches(7.1), Inches(10), Inches(0.4))
    footer_rect4.fill.solid()
    footer_rect4.fill.fore_color.rgb = footer_blue
    footer_rect4.line.fill.background()

    prs.save('NEUROOPS_Presentation.pptx')
    print("Presentation created successfully as NEUROOPS_Presentation.pptx")

if __name__ == "__main__":
    create_presentation()
